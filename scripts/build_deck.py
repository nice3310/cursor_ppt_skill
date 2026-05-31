"""Build a .pptx from an outline.yaml.

Flow:
  1. Load + validate the outline (hard schema only; soft warnings run separately).
  2. Resolve template path: meta.template -> templates/default.pptx -> python-pptx default.
  3. Open template, REMOVE any pre-existing slides (we want a clean deck).
  4. Pick a theme (meta.theme) and render each slide with the designed renderers.
  5. Save to meta.output.

Usage:
    python scripts/build_deck.py path/to/outline.yaml
    python scripts/build_deck.py path/to/outline.yaml --no-mermaid
    python scripts/build_deck.py path/to/outline.yaml -o custom.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pydantic import ValidationError

sys.path.insert(0, str(Path(__file__).parent))
from lib.renderers import render_slide  # noqa: E402
from lib.schema import Outline  # noqa: E402
from lib.theme import Theme, get_theme  # noqa: E402


SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
DEFAULT_TEMPLATE = SKILL_DIR / "templates" / "default.pptx"


def _hex(color) -> str:
    return str(color)


def _load(path: Path) -> Outline:
    data = yaml.safe_load(path.read_text(encoding="utf-8"))
    return Outline.model_validate(data)


def _resolve_template(outline: Outline, outline_dir: Path) -> Path | None:
    if outline.meta.template:
        p = Path(outline.meta.template)
        resolved = p if p.is_absolute() else (outline_dir / p).resolve()
        if not resolved.exists():
            raise FileNotFoundError(f"meta.template not found: {resolved}")
        return resolved
    if DEFAULT_TEMPLATE.exists():
        return DEFAULT_TEMPLATE
    return None


def _wipe_slides(prs) -> None:
    sldIdLst = prs.slides._sldIdLst
    for sld in list(sldIdLst):
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sld)


def _make_mermaid_renderer(enabled: bool, theme: Theme):
    if not enabled:
        return None
    from render_mermaid import MermaidRenderError, default_theme_vars, render  # type: ignore

    theme_vars = default_theme_vars(
        accent=_hex(theme.accent),
        accent2=_hex(theme.accent2),
        ink=_hex(theme.ink),
        line=_hex(theme.ink_soft),
        surface=_hex(theme.surface),
        font=f"{theme.font_body}, Segoe UI, Helvetica, Arial, sans-serif",
    )

    def _do(source: str, out_path: Path) -> Path:
        return render(source, out_path, theme_vars=theme_vars)

    return _do


def build(outline_path: Path, output_override: Path | None, *, mermaid: bool) -> Path:
    outline = _load(outline_path)
    outline_dir = outline_path.resolve().parent

    template = _resolve_template(outline, outline_dir)
    prs = Presentation(str(template)) if template else Presentation()
    _wipe_slides(prs)

    theme = get_theme(outline.meta.theme)
    mermaid_renderer = _make_mermaid_renderer(mermaid, theme)

    footer = outline.meta.footer or outline.meta.title
    total = len(outline.slides)
    section_counter = 0

    for i, slide_model in enumerate(outline.slides, start=1):
        if slide_model.type == "section":
            section_counter += 1
        render_slide(
            prs,
            slide_model,
            theme,
            base_dir=outline_dir,
            mermaid_renderer=mermaid_renderer,
            footer=footer,
            page=i,
            total=total,
            section_index=section_counter if slide_model.type == "section" else None,
        )

    out_path = output_override or (
        Path(outline.meta.output)
        if Path(outline.meta.output).is_absolute()
        else outline_dir / outline.meta.output
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build a .pptx from outline.yaml")
    parser.add_argument("outline", type=Path)
    parser.add_argument("-o", "--output", type=Path, help="Override meta.output")
    parser.add_argument(
        "--no-mermaid",
        action="store_true",
        help="Skip mermaid rendering (diagrams become placeholders)",
    )
    args = parser.parse_args(argv)

    if not args.outline.exists():
        print(f"error: {args.outline} not found", file=sys.stderr)
        return 2

    try:
        out = build(args.outline, args.output, mermaid=not args.no_mermaid)
    except ValidationError as e:
        print("schema errors (run validate_outline.py for details):", file=sys.stderr)
        for err in e.errors():
            loc = ".".join(str(p) for p in err["loc"])
            print(f"  {loc}: {err['msg']}", file=sys.stderr)
        return 1
    except FileNotFoundError as e:
        print(f"error: {e}", file=sys.stderr)
        return 1

    print(f"wrote {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
