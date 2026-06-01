"""Per-slide-type renderers with a real design system.

Each renderer draws its own shapes on a chosen layout so we inherit the
template's master while controlling the visual treatment: kicker eyebrows,
accent bars, surface cards, full-bleed impact slides, a dark syntax-
highlighted code card, zebra tables, and a consistent footer.

Positioning assumes a 16:9 slide (13.333 x 7.5 in) but reads the real size
from the presentation, so it degrades reasonably on 4:3.

Renderers now accept a `Style` parameter that controls *layout geometry
and decorative language* independently of color (which is still `Theme`).
"""

from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationT
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu, Inches, Pt

from .code_highlight import highlight_lines, style_background_hex
from .layouts import pick_blank_layout
from .schema import (
    BulletsSlide,
    CodeSlide,
    ComparisonSlide,
    DiagramSlide,
    ImageSlide,
    QASlide,
    QuoteSlide,
    SectionSlide,
    Slide,
    StatementSlide,
    SummarySlide,
    TableSlide,
    TitleSlide,
)
from .styles import Style
from .theme import Theme

try:
    from PIL import Image
except Exception:  # noqa: BLE001
    Image = None


# Vertical layout constants (inches) for content slides.
KICKER_TOP = 0.55
TITLE_TOP = 0.92
BAR_TOP = 1.64
HEADLINE_TOP = 1.78
CONTENT_TOP = 2.45


# ---------------------------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------------------------


def _slide_size_in(prs: PresentationT) -> tuple[float, float]:
    return Emu(prs.slide_width).inches, Emu(prs.slide_height).inches


def _new_slide(prs: PresentationT) -> PptxSlide:
    return prs.slides.add_slide(pick_blank_layout(prs))


def _bg(slide: PptxSlide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, shape_type, left, top, width, height, fill=None, line=None):
    shp = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def _soft_shadow(shape, blur_in=0.06, dist_in=0.035, alpha=24, dir_deg=90):
    """Inject a subtle outer drop shadow so cards/images read as layered."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        return
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(
        qn("a:outerShdw"),
        {
            "blurRad": str(int(Inches(blur_in))),
            "dist": str(int(Inches(dist_in))),
            "dir": str(int(dir_deg * 60000)),
            "rotWithShape": "0",
        },
    )
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def _textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return box, tf


def _run(p, text, *, size, color, bold=False, italic=False, font=None, tracking=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if font:
        r.font.name = font
    if tracking is not None:
        r._r.get_or_add_rPr().set("spc", str(int(tracking * 100)))
    return r


def _simple_text(
    slide,
    text,
    *,
    left,
    top,
    width,
    height,
    size,
    color,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=None,
    tracking=None,
    line_spacing=None,
):
    box, tf = _textbox(slide, left=left, top=top, width=width, height=height, anchor=anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    _run(p, text, size=size, color=color, bold=bold, italic=italic, font=font, tracking=tracking)
    return box


# ---------------------------------------------------------------------------
# effective geometry helpers (style-aware)
# ---------------------------------------------------------------------------

def _eff_margin(theme: Theme, style: Style) -> float:
    return theme.margin_left * style.margin_scale


def _eff_content_top(style: Style) -> float:
    return CONTENT_TOP + style.content_top_offset


# ---------------------------------------------------------------------------
# shared chrome (kicker / title / headline / footer) — style-aware
# ---------------------------------------------------------------------------


def _kicker(slide, text, theme: Theme, style: Style, slide_w):
    ml = _eff_margin(theme, style)
    _simple_text(
        slide,
        text.upper(),
        left=ml,
        top=KICKER_TOP,
        width=slide_w - 2 * ml,
        height=0.3,
        size=theme.size_kicker,
        color=theme.accent,
        bold=True,
        font=theme.font_display,
        tracking=2.2,
    )
    if style.use_kicker_line:
        _rect(slide, MSO_SHAPE.RECTANGLE, ml, KICKER_TOP + 0.28, 1.2, 0.015, fill=theme.accent)


def _title_chrome(slide, text, theme: Theme, style: Style, slide_w):
    """Draw title text + accent bar based on style."""
    ml = _eff_margin(theme, style)
    _simple_text(
        slide,
        text,
        left=ml,
        top=TITLE_TOP,
        width=slide_w - 2 * ml,
        height=0.7,
        size=theme.size_title,
        color=theme.ink,
        bold=True,
        font=theme.font_display,
    )
    # Accent bar variants
    bar = style.accent_bar
    if bar == "short_bar":
        _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, ml, BAR_TOP, theme.accent_bar_w, theme.accent_bar_h, fill=theme.accent)
    elif bar == "underline":
        _rect(slide, MSO_SHAPE.RECTANGLE, ml, BAR_TOP, slide_w - 2 * ml, 0.03, fill=theme.accent)
    elif bar == "left_thick":
        _rect(slide, MSO_SHAPE.RECTANGLE, ml - 0.18, TITLE_TOP, 0.08, 0.65, fill=theme.accent)
    elif bar == "dot":
        _rect(slide, MSO_SHAPE.OVAL, ml, BAR_TOP, 0.14, 0.14, fill=theme.accent)
    # "none" — no accent bar


def _headline(slide, text, theme: Theme, style: Style, slide_w):
    ml = _eff_margin(theme, style)
    _simple_text(
        slide,
        text,
        left=ml,
        top=HEADLINE_TOP,
        width=slide_w - 2 * ml,
        height=0.6,
        size=theme.size_headline,
        color=theme.ink_soft,
        font=theme.font_body,
        italic=style.headline_italic,
        line_spacing=1.05,
    )


def _footer(slide, theme: Theme, style: Style, slide_w, slide_h, footer_text, page, total):
    if not footer_text and not page:
        return
    ml = _eff_margin(theme, style)
    y = slide_h - theme.margin_bottom - 0.05
    if style.footer_style == "rule":
        _rect(slide, MSO_SHAPE.RECTANGLE, ml, y, slide_w - 2 * ml, 0.012, fill=theme.rule)
    if footer_text:
        _simple_text(
            slide,
            footer_text,
            left=ml,
            top=y + 0.04,
            width=(slide_w - 2 * ml) * 0.7,
            height=0.3,
            size=theme.size_footer,
            color=theme.ink_soft,
            font=theme.font_body,
        )
    if page:
        _simple_text(
            slide,
            f"{page} / {total}" if total else str(page),
            left=slide_w - ml - 1.5,
            top=y + 0.04,
            width=1.5,
            height=0.3,
            size=theme.size_footer,
            color=theme.ink_soft,
            align=PP_ALIGN.RIGHT,
            font=theme.font_body,
        )


def _set_notes(slide, text: Optional[str], source: Optional[str] = None) -> None:
    parts = []
    if text:
        parts.append(text)
    if source:
        parts.append(f"Image source: {source}")
    if not parts:
        return
    slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)


def _content_chrome(slide, model, theme, style, slide_w, slide_h, footer, page, total):
    """Draw background + kicker + title + accent bar + headline + footer."""
    _bg(slide, theme.background)
    kicker = getattr(model, "kicker", None)
    if kicker:
        _kicker(slide, kicker, theme, style, slide_w)
    _title_chrome(slide, model.title, theme, style, slide_w)
    if getattr(model, "headline", None):
        _headline(slide, model.headline, theme, style, slide_w)
    _footer(slide, theme, style, slide_w, slide_h, footer, page, total)


# ---------------------------------------------------------------------------
# bullet rows — style-aware markers
# ---------------------------------------------------------------------------


def _bullet_rows(
    slide,
    items,
    *,
    left,
    top,
    width,
    height,
    theme: Theme,
    style: Style,
    numbered=False,
    text_size=None,
    marker_color=None,
):
    text_size = text_size or theme.size_body
    marker_color = marker_color or theme.accent
    n = len(items)
    row_h = min(0.95, height / n)
    marker_style = style.bullet_marker

    for i, item in enumerate(items):
        ry = top + i * row_h

        if numbered:
            # Numbered chips always use rounded rect
            chip = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, ry + 0.04, 0.34, 0.34, fill=marker_color)
            _, ctf = chip.text_frame, chip.text_frame
            ctf.word_wrap = False
            ctf.margin_left = Emu(0)
            ctf.margin_right = Emu(0)
            ctf.margin_top = Emu(0)
            ctf.margin_bottom = Emu(0)
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            _run(cp, str(i + 1), size=theme.size_caption, color=theme.on_accent, bold=True, font=theme.font_display)
            text_left = left + 0.5
        else:
            marker_sz = text_size / 72.0
            my = ry + marker_sz / 2.0

            if marker_style == "rounded_rect":
                _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, my, 0.16, 0.16, fill=marker_color)
                text_left = left + 0.38
            elif marker_style == "circle":
                _rect(slide, MSO_SHAPE.OVAL, left + 0.02, my + 0.02, 0.12, 0.12, fill=marker_color)
                text_left = left + 0.34
            elif marker_style == "dash":
                _rect(slide, MSO_SHAPE.RECTANGLE, left, my + 0.06, 0.22, 0.04, fill=marker_color)
                text_left = left + 0.38
            elif marker_style == "pill":
                _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, my + 0.01, 0.24, 0.14, fill=marker_color)
                text_left = left + 0.40
            elif marker_style == "diamond":
                # Diamond = rotated square — approximate with a small shape
                _rect(slide, MSO_SHAPE.DIAMOND, left + 0.01, my, 0.15, 0.15, fill=marker_color)
                text_left = left + 0.36
            else:
                _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, my, 0.16, 0.16, fill=marker_color)
                text_left = left + 0.38

        _simple_text(
            slide,
            item,
            left=text_left,
            top=ry,
            width=width - (text_left - left),
            height=row_h,
            size=text_size,
            color=theme.ink,
            font=theme.font_body,
            anchor=MSO_ANCHOR.TOP,
            line_spacing=1.02,
        )


# ---------------------------------------------------------------------------
# image fitting
# ---------------------------------------------------------------------------


def _image_aspect(path: Path) -> Optional[float]:
    if Image is None:
        return None
    try:
        with Image.open(path) as im:
            w, h = im.size
            return w / h if h else None
    except Exception:  # noqa: BLE001
        return None


def _place_image_contain(slide, path: Path, box, theme: Theme, shadow=True):
    """Place an image fully inside box (inches), preserving aspect ratio, centered."""
    bx, by, bw, bh = box
    ar = _image_aspect(path)
    if ar is None:
        # Unknown size: fill the box and let it be.
        pic = slide.shapes.add_picture(str(path), Inches(bx), Inches(by), Inches(bw), Inches(bh))
    else:
        box_ar = bw / bh
        if ar > box_ar:
            w = bw
            h = bw / ar
        else:
            h = bh
            w = bh * ar
        x = bx + (bw - w) / 2
        y = by + (bh - h) / 2
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.line.color.rgb = theme.surface_edge
    pic.line.width = Pt(0.75)
    if shadow:
        _soft_shadow(pic)
    return pic


def _placeholder(slide, box, theme, text, color=None):
    bx, by, bw, bh = box
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, fill=theme.surface, line=theme.surface_edge)
    _simple_text(
        slide,
        text,
        left=bx,
        top=by,
        width=bw,
        height=bh,
        size=theme.size_body,
        color=color or theme.ink_soft,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        font=theme.font_body,
    )


# ===================================================================
# TITLE SLIDE VARIANTS
# ===================================================================

def _title_rail(slide, s, theme, style, slide_w, slide_h):
    """Editorial: left accent rail + text offset right."""
    _bg(slide, theme.background)
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.28, slide_h, fill=theme.accent)
    ml = _eff_margin(theme, style)

    _simple_text(
        slide, s.title,
        left=ml + 0.3, top=slide_h * 0.30,
        width=slide_w - ml - theme.margin_right - 0.3, height=1.8,
        size=theme.size_section, color=theme.ink, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )
    if s.subtitle:
        _simple_text(
            slide, s.subtitle,
            left=ml + 0.32, top=slide_h * 0.30 + 1.85,
            width=slide_w - ml - theme.margin_right - 0.3, height=0.5,
            size=theme.size_headline, color=theme.ink_soft, font=theme.font_body,
        )
    meta_bits = [b for b in (s.presenter, str(s.date) if s.date else None) if b]
    if meta_bits:
        ry = slide_h - theme.margin_bottom - 0.5
        _rect(slide, MSO_SHAPE.RECTANGLE, ml + 0.32, ry, 0.5, 0.03, fill=theme.accent)
        _simple_text(
            slide, "   \u00b7   ".join(meta_bits),
            left=ml + 0.32, top=ry + 0.08,
            width=slide_w - ml - theme.margin_right, height=0.4,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


def _title_centered(slide, s, theme, style, slide_w, slide_h):
    """Minimal: centered, clean, no rail."""
    _bg(slide, theme.background)

    _simple_text(
        slide, s.title,
        left=1.5, top=slide_h * 0.32,
        width=slide_w - 3.0, height=1.8,
        size=theme.size_section + 2, color=theme.ink, bold=True, font=theme.font_display,
        line_spacing=1.0, align=PP_ALIGN.CENTER,
    )
    # Thin line under title
    cx = slide_w / 2
    _rect(slide, MSO_SHAPE.RECTANGLE, cx - 1.0, slide_h * 0.32 + 1.9, 2.0, 0.02, fill=theme.rule)

    if s.subtitle:
        _simple_text(
            slide, s.subtitle,
            left=1.5, top=slide_h * 0.32 + 2.1,
            width=slide_w - 3.0, height=0.5,
            size=theme.size_headline, color=theme.ink_soft, font=theme.font_body,
            align=PP_ALIGN.CENTER, italic=True,
        )
    meta_bits = [b for b in (s.presenter, str(s.date) if s.date else None) if b]
    if meta_bits:
        _simple_text(
            slide, "   \u00b7   ".join(meta_bits),
            left=1.5, top=slide_h - 1.2,
            width=slide_w - 3.0, height=0.4,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
            align=PP_ALIGN.CENTER,
        )


def _title_full_bleed(slide, s, theme, style, slide_w, slide_h):
    """Corporate: full-width accent band at top."""
    _bg(slide, theme.background)
    # Top accent band
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h * 0.42, fill=theme.accent)
    ml = _eff_margin(theme, style)

    _simple_text(
        slide, s.title,
        left=ml, top=slide_h * 0.10,
        width=slide_w - 2 * ml, height=1.6,
        size=theme.size_section + 4, color=theme.on_accent, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )
    if s.subtitle:
        _simple_text(
            slide, s.subtitle,
            left=ml, top=slide_h * 0.46,
            width=slide_w - 2 * ml, height=0.5,
            size=theme.size_headline + 2, color=theme.ink_soft, font=theme.font_body,
        )
    meta_bits = [b for b in (s.presenter, str(s.date) if s.date else None) if b]
    if meta_bits:
        ry = slide_h - theme.margin_bottom - 0.5
        _simple_text(
            slide, "   \u00b7   ".join(meta_bits),
            left=ml, top=ry,
            width=slide_w - 2 * ml, height=0.4,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


def _title_split(slide, s, theme, style, slide_w, slide_h):
    """Modern: left half accent, right half content."""
    _bg(slide, theme.background)
    # Left half accent
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_w * 0.42, slide_h, fill=theme.accent)

    right_x = slide_w * 0.48
    right_w = slide_w - right_x - theme.margin_right

    _simple_text(
        slide, s.title,
        left=right_x, top=slide_h * 0.25,
        width=right_w, height=1.8,
        size=theme.size_section, color=theme.ink, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )
    if s.subtitle:
        _simple_text(
            slide, s.subtitle,
            left=right_x, top=slide_h * 0.25 + 1.85,
            width=right_w, height=0.5,
            size=theme.size_headline, color=theme.ink_soft, font=theme.font_body,
        )
    meta_bits = [b for b in (s.presenter, str(s.date) if s.date else None) if b]
    if meta_bits:
        _simple_text(
            slide, "   \u00b7   ".join(meta_bits),
            left=right_x, top=slide_h - 1.0,
            width=right_w, height=0.4,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


def _title_diagonal(slide, s, theme, style, slide_w, slide_h):
    """Vibrant: diagonal accent element."""
    _bg(slide, theme.background)
    # Diagonal created via a large accent shape in the corner
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_w * 0.35, slide_h, fill=theme.accent)
    _rect(slide, MSO_SHAPE.RECTANGLE, slide_w * 0.35, 0, slide_w * 0.08, slide_h, fill=theme.accent2)
    ml = _eff_margin(theme, style)

    _simple_text(
        slide, s.title,
        left=slide_w * 0.48, top=slide_h * 0.25,
        width=slide_w * 0.48, height=1.8,
        size=theme.size_section, color=theme.ink, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )
    if s.subtitle:
        _simple_text(
            slide, s.subtitle,
            left=slide_w * 0.48, top=slide_h * 0.25 + 1.85,
            width=slide_w * 0.48, height=0.5,
            size=theme.size_headline, color=theme.ink_soft, font=theme.font_body,
        )
    meta_bits = [b for b in (s.presenter, str(s.date) if s.date else None) if b]
    if meta_bits:
        _simple_text(
            slide, "   \u00b7   ".join(meta_bits),
            left=slide_w * 0.48, top=slide_h - 1.0,
            width=slide_w * 0.48, height=0.4,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


_TITLE_VARIANTS = {
    "rail": _title_rail,
    "centered": _title_centered,
    "full_bleed": _title_full_bleed,
    "split": _title_split,
    "diagonal": _title_diagonal,
}


# ===================================================================
# SECTION SLIDE VARIANTS
# ===================================================================

def _section_dark_full(slide, s, theme, style, slide_w, slide_h, num):
    """Editorial: full dark background + ghosted number."""
    _bg(slide, theme.ink)
    ml = _eff_margin(theme, style)
    if num:
        _simple_text(
            slide, num,
            left=ml, top=slide_h * 0.18,
            width=6.0, height=2.6,
            size=theme.size_section_num, color=theme.ink_soft, bold=True, font=theme.font_display,
        )
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, ml + 0.04, slide_h * 0.60, 0.8, 0.09, fill=theme.accent)
    _simple_text(
        slide, s.title,
        left=ml, top=slide_h * 0.62,
        width=slide_w - 2 * ml, height=1.6,
        size=theme.size_section, color=theme.on_accent, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )


def _section_accent_band(slide, s, theme, style, slide_w, slide_h, num):
    """Corporate: accent band across middle."""
    _bg(slide, theme.background)
    band_top = slide_h * 0.35
    band_h = slide_h * 0.30
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, band_top, slide_w, band_h, fill=theme.accent)
    ml = _eff_margin(theme, style)
    if num:
        _simple_text(
            slide, num,
            left=ml, top=band_top + 0.15,
            width=1.5, height=band_h - 0.3,
            size=theme.size_section_num * 0.6, color=theme.on_accent_soft, bold=True, font=theme.font_display,
        )
    _simple_text(
        slide, s.title,
        left=ml + (1.8 if num else 0), top=band_top,
        width=slide_w - 2 * ml - (1.8 if num else 0), height=band_h,
        size=theme.size_section, color=theme.on_accent, bold=True, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _section_clean(slide, s, theme, style, slide_w, slide_h, num):
    """Minimal: white background, large light number."""
    _bg(slide, theme.background)
    ml = _eff_margin(theme, style)
    if num:
        _simple_text(
            slide, num,
            left=ml, top=slide_h * 0.15,
            width=6.0, height=2.6,
            size=theme.size_section_num, color=theme.surface, bold=True, font=theme.font_display,
        )
    _simple_text(
        slide, s.title,
        left=ml, top=slide_h * 0.45,
        width=slide_w - 2 * ml, height=1.6,
        size=theme.size_section, color=theme.ink, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )
    # Thin line under
    _rect(slide, MSO_SHAPE.RECTANGLE, ml, slide_h * 0.45 + 1.5, 2.0, 0.02, fill=theme.accent)


def _section_gradient(slide, s, theme, style, slide_w, slide_h, num):
    """Modern: simulate gradient via layered accent tones."""
    _bg(slide, theme.accent)
    # Overlay a softer band
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, slide_h * 0.6, slide_w, slide_h * 0.4, fill=theme.accent_soft)
    ml = _eff_margin(theme, style)
    if num:
        _simple_text(
            slide, num,
            left=ml, top=slide_h * 0.12,
            width=6.0, height=2.0,
            size=theme.size_section_num * 0.7, color=theme.on_accent_soft, bold=True, font=theme.font_display,
        )
    _simple_text(
        slide, s.title,
        left=ml, top=slide_h * 0.35,
        width=slide_w - 2 * ml, height=1.6,
        size=theme.size_section, color=theme.on_accent, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )


def _section_color_block(slide, s, theme, style, slide_w, slide_h, num):
    """Vibrant: full accent2 background."""
    _bg(slide, theme.accent2)
    ml = _eff_margin(theme, style)
    if num:
        _simple_text(
            slide, num,
            left=ml, top=slide_h * 0.15,
            width=6.0, height=2.6,
            size=theme.size_section_num, color=theme.on_accent_soft, bold=True, font=theme.font_display,
        )
    _simple_text(
        slide, s.title,
        left=ml, top=slide_h * 0.55,
        width=slide_w - 2 * ml, height=1.6,
        size=theme.size_section, color=theme.on_accent, bold=True, font=theme.font_display,
        line_spacing=1.0,
    )


_SECTION_VARIANTS = {
    "dark_full": _section_dark_full,
    "accent_band": _section_accent_band,
    "clean": _section_clean,
    "gradient": _section_gradient,
    "color_block": _section_color_block,
}


# ===================================================================
# STATEMENT SLIDE VARIANTS
# ===================================================================

def _statement_full_accent(slide, s, theme, style, slide_w, slide_h):
    """Editorial: full accent background."""
    _bg(slide, theme.accent)
    if s.kicker:
        _simple_text(
            slide, s.kicker.upper(),
            left=theme.margin_left + 0.3, top=slide_h * 0.26,
            width=slide_w - 2 * theme.margin_left - 0.6, height=0.4,
            size=theme.size_kicker, color=theme.on_accent_soft, bold=True,
            font=theme.font_display, tracking=2.4, align=PP_ALIGN.CENTER,
        )
    _simple_text(
        slide, s.headline,
        left=theme.margin_left + 0.3, top=slide_h * 0.30,
        width=slide_w - 2 * theme.margin_left - 0.6, height=slide_h * 0.42,
        size=theme.size_statement, color=theme.on_accent, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        line_spacing=1.05,
    )


def _statement_centered_large(slide, s, theme, style, slide_w, slide_h):
    """Corporate: white bg, extra-large text."""
    _bg(slide, theme.background)
    if s.kicker:
        _simple_text(
            slide, s.kicker.upper(),
            left=1.0, top=slide_h * 0.22,
            width=slide_w - 2.0, height=0.4,
            size=theme.size_kicker, color=theme.accent, bold=True,
            font=theme.font_display, tracking=2.4, align=PP_ALIGN.CENTER,
        )
    _simple_text(
        slide, s.headline,
        left=1.0, top=slide_h * 0.28,
        width=slide_w - 2.0, height=slide_h * 0.48,
        size=theme.size_statement + 6, color=theme.ink, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        line_spacing=1.1,
    )
    # Accent dot below
    _rect(slide, MSO_SHAPE.OVAL, slide_w / 2 - 0.08, slide_h * 0.78, 0.16, 0.16, fill=theme.accent)


def _statement_with_line(slide, s, theme, style, slide_w, slide_h):
    """Minimal: centered text + thin lines above/below."""
    _bg(slide, theme.background)
    cx = slide_w / 2
    if s.kicker:
        _simple_text(
            slide, s.kicker.upper(),
            left=1.0, top=slide_h * 0.24,
            width=slide_w - 2.0, height=0.4,
            size=theme.size_kicker, color=theme.accent, bold=True,
            font=theme.font_display, tracking=2.4, align=PP_ALIGN.CENTER,
        )
    _rect(slide, MSO_SHAPE.RECTANGLE, cx - 2.0, slide_h * 0.34, 4.0, 0.015, fill=theme.rule)
    _simple_text(
        slide, s.headline,
        left=1.5, top=slide_h * 0.36,
        width=slide_w - 3.0, height=slide_h * 0.32,
        size=theme.size_statement, color=theme.ink, bold=False, italic=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        line_spacing=1.1,
    )
    _rect(slide, MSO_SHAPE.RECTANGLE, cx - 2.0, slide_h * 0.70, 4.0, 0.015, fill=theme.rule)


def _statement_gradient_bg(slide, s, theme, style, slide_w, slide_h):
    """Modern: simulate gradient with layered fills."""
    _bg(slide, theme.accent)
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, slide_h * 0.55, slide_w, slide_h * 0.45, fill=theme.accent_soft)
    if s.kicker:
        _simple_text(
            slide, s.kicker.upper(),
            left=1.0, top=slide_h * 0.22,
            width=slide_w - 2.0, height=0.4,
            size=theme.size_kicker, color=theme.on_accent_soft, bold=True,
            font=theme.font_display, tracking=2.4, align=PP_ALIGN.CENTER,
        )
    _simple_text(
        slide, s.headline,
        left=1.2, top=slide_h * 0.28,
        width=slide_w - 2.4, height=slide_h * 0.42,
        size=theme.size_statement, color=theme.on_accent, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        line_spacing=1.05,
    )


def _statement_dual_color(slide, s, theme, style, slide_w, slide_h):
    """Vibrant: two-tone background."""
    _bg(slide, theme.accent)
    _rect(slide, MSO_SHAPE.RECTANGLE, slide_w * 0.55, 0, slide_w * 0.45, slide_h, fill=theme.accent2)
    if s.kicker:
        _simple_text(
            slide, s.kicker.upper(),
            left=1.0, top=slide_h * 0.26,
            width=slide_w - 2.0, height=0.4,
            size=theme.size_kicker, color=theme.on_accent_soft, bold=True,
            font=theme.font_display, tracking=2.4, align=PP_ALIGN.CENTER,
        )
    _simple_text(
        slide, s.headline,
        left=1.2, top=slide_h * 0.30,
        width=slide_w - 2.4, height=slide_h * 0.44,
        size=theme.size_statement, color=theme.on_accent, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        line_spacing=1.05,
    )


_STATEMENT_VARIANTS = {
    "full_accent": _statement_full_accent,
    "centered_large": _statement_centered_large,
    "with_line": _statement_with_line,
    "gradient_bg": _statement_gradient_bg,
    "dual_color": _statement_dual_color,
}


# ===================================================================
# QUOTE SLIDE VARIANTS
# ===================================================================

def _quote_big_mark(slide, s, theme, style, slide_w, slide_h):
    """Editorial: large quotation mark + surface bg."""
    _bg(slide, theme.surface)
    _simple_text(
        slide, "\u201C",
        left=theme.margin_left, top=slide_h * 0.10,
        width=2.0, height=2.0,
        size=theme.size_quote_mark, color=theme.accent, bold=True, font=theme.font_display,
    )
    _simple_text(
        slide, s.quote,
        left=theme.margin_left + 0.5, top=slide_h * 0.34,
        width=slide_w - 2 * theme.margin_left - 1.0, height=slide_h * 0.34,
        size=theme.size_quote, color=theme.ink, bold=False, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1,
    )
    if s.attribution:
        ry = slide_h * 0.74
        _rect(slide, MSO_SHAPE.RECTANGLE, theme.margin_left + 0.55, ry, 0.5, 0.03, fill=theme.accent)
        _simple_text(
            slide, s.attribution,
            left=theme.margin_left + 0.55, top=ry + 0.1,
            width=slide_w - 2 * theme.margin_left, height=0.5,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


def _quote_left_bar(slide, s, theme, style, slide_w, slide_h):
    """Corporate: thick left accent bar."""
    _bg(slide, theme.background)
    _rect(slide, MSO_SHAPE.RECTANGLE, theme.margin_left, slide_h * 0.25, 0.12, slide_h * 0.45, fill=theme.accent)
    _simple_text(
        slide, s.quote,
        left=theme.margin_left + 0.5, top=slide_h * 0.28,
        width=slide_w - theme.margin_left - theme.margin_right - 0.8, height=slide_h * 0.36,
        size=theme.size_quote, color=theme.ink, bold=False, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1,
    )
    if s.attribution:
        _simple_text(
            slide, f"\u2014 {s.attribution}",
            left=theme.margin_left + 0.5, top=slide_h * 0.70,
            width=slide_w - 2 * theme.margin_left, height=0.5,
            size=theme.size_attribution, color=theme.ink_soft, bold=True, font=theme.font_body,
        )


def _quote_italic_center(slide, s, theme, style, slide_w, slide_h):
    """Minimal: centered italic, no decorations."""
    _bg(slide, theme.background)
    _simple_text(
        slide, s.quote,
        left=2.0, top=slide_h * 0.30,
        width=slide_w - 4.0, height=slide_h * 0.35,
        size=theme.size_quote - 2, color=theme.ink, bold=False, italic=True, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line_spacing=1.15,
    )
    if s.attribution:
        _rect(slide, MSO_SHAPE.RECTANGLE, slide_w / 2 - 0.5, slide_h * 0.70, 1.0, 0.02, fill=theme.rule)
        _simple_text(
            slide, s.attribution,
            left=2.0, top=slide_h * 0.73,
            width=slide_w - 4.0, height=0.5,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
            align=PP_ALIGN.CENTER,
        )


def _quote_card(slide, s, theme, style, slide_w, slide_h):
    """Modern: quote inside a card with shadow."""
    _bg(slide, theme.background)
    card_x = theme.margin_left + 0.5
    card_w = slide_w - 2 * (theme.margin_left + 0.5)
    card_y = slide_h * 0.20
    card_h = slide_h * 0.55
    card = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h, fill=theme.surface)
    _soft_shadow(card, blur_in=0.1, dist_in=0.06, alpha=30)
    _simple_text(
        slide, "\u201C",
        left=card_x + 0.3, top=card_y + 0.2,
        width=1.0, height=1.0,
        size=80, color=theme.accent, bold=True, font=theme.font_display,
    )
    _simple_text(
        slide, s.quote,
        left=card_x + 0.6, top=card_y + card_h * 0.25,
        width=card_w - 1.2, height=card_h * 0.50,
        size=theme.size_quote - 4, color=theme.ink, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1,
    )
    if s.attribution:
        _simple_text(
            slide, f"\u2014 {s.attribution}",
            left=card_x + 0.6, top=card_y + card_h - 0.65,
            width=card_w - 1.2, height=0.5,
            size=theme.size_attribution, color=theme.ink_soft, font=theme.font_body,
        )


def _quote_accent_bg(slide, s, theme, style, slide_w, slide_h):
    """Vibrant: accent background + white text."""
    _bg(slide, theme.accent)
    _simple_text(
        slide, "\u201C",
        left=theme.margin_left, top=slide_h * 0.10,
        width=2.0, height=1.5,
        size=100, color=theme.on_accent_soft, bold=True, font=theme.font_display,
    )
    _simple_text(
        slide, s.quote,
        left=theme.margin_left + 0.5, top=slide_h * 0.30,
        width=slide_w - 2 * theme.margin_left - 1.0, height=slide_h * 0.38,
        size=theme.size_quote, color=theme.on_accent, font=theme.font_display,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1,
    )
    if s.attribution:
        ry = slide_h * 0.74
        _rect(slide, MSO_SHAPE.RECTANGLE, theme.margin_left + 0.55, ry, 0.5, 0.03, fill=theme.on_accent)
        _simple_text(
            slide, s.attribution,
            left=theme.margin_left + 0.55, top=ry + 0.1,
            width=slide_w - 2 * theme.margin_left, height=0.5,
            size=theme.size_attribution, color=theme.on_accent_soft, font=theme.font_body,
        )


_QUOTE_VARIANTS = {
    "big_mark": _quote_big_mark,
    "left_bar": _quote_left_bar,
    "italic_center": _quote_italic_center,
    "card": _quote_card,
    "accent_bg": _quote_accent_bg,
}


# ===================================================================
# CARD DRAWING HELPERS (for comparison / code)
# ===================================================================

def _draw_card(slide, left, top, width, height, theme, style, fill_override=None):
    """Draw a card shape according to the style's card_style."""
    cs = style.card_style
    fill = fill_override

    if cs == "bordered_round":
        card = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill=fill or theme.surface, line=theme.surface_edge)
        _soft_shadow(card, alpha=14)
        return card
    elif cs == "flat_fill":
        return _rect(slide, MSO_SHAPE.RECTANGLE, left, top, width, height,
                      fill=fill or theme.surface)
    elif cs == "border_only":
        if fill is not None:
            return _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=theme.surface_edge)
        return _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      line=theme.surface_edge)
    elif cs == "heavy_shadow":
        card = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill=fill or theme.surface)
        _soft_shadow(card, blur_in=0.12, dist_in=0.06, alpha=35)
        return card
    elif cs == "left_accent":
        card = _rect(slide, MSO_SHAPE.RECTANGLE, left, top, width, height,
                      fill=fill or theme.surface)
        _rect(slide, MSO_SHAPE.RECTANGLE, left, top, 0.08, height, fill=theme.accent)
        return card
    else:
        card = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill=fill or theme.surface, line=theme.surface_edge)
        _soft_shadow(card, alpha=14)
        return card


# ===================================================================
# RENDERERS (main entry points)
# ===================================================================


def render_title(prs, s: TitleSlide, theme: Theme, style: Style, *, base_dir, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    fn = _TITLE_VARIANTS.get(style.title_variant, _title_rail)
    fn(slide, s, theme, style, slide_w, slide_h)
    _set_notes(slide, s.notes)
    return slide


def render_section(prs, s: SectionSlide, theme: Theme, style: Style, *, section_index=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    num = s.number or (f"{section_index:02d}" if section_index is not None else None)
    fn = _SECTION_VARIANTS.get(style.section_variant, _section_dark_full)
    fn(slide, s, theme, style, slide_w, slide_h, num)
    _set_notes(slide, s.notes)
    return slide


def render_statement(prs, s: StatementSlide, theme: Theme, style: Style, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    fn = _STATEMENT_VARIANTS.get(style.statement_variant, _statement_full_accent)
    fn(slide, s, theme, style, slide_w, slide_h)
    _set_notes(slide, s.notes)
    return slide


# ===================================================================
# MACRO-LAYOUT VARIANTS: BULLETS
# ===================================================================

def _bullets_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)
    _bullet_rows(
        slide, s.bullets,
        left=ml, top=ct, width=slide_w - 2 * ml, height=slide_h - ct - theme.margin_bottom - 0.3,
        theme=theme, style=style, numbered=s.numbered,
    )

def _bullets_magazine(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _bg(slide, theme.background)
    _footer(slide, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    left_w = (slide_w - 2 * ml) * 0.4
    right_w = (slide_w - 2 * ml) * 0.55
    right_x = slide_w - ml - right_w
    
    if getattr(s, "kicker", None):
        _simple_text(slide, s.kicker.upper(), left=ml, top=slide_h * 0.25, width=left_w, height=0.4, size=theme.size_kicker, color=theme.accent, bold=True, font=theme.font_display, tracking=2.4)
    _simple_text(slide, s.title, left=ml, top=slide_h * 0.3, width=left_w, height=1.6, size=theme.size_title + 8, color=theme.ink, bold=True, font=theme.font_display, line_spacing=0.9)
    if getattr(s, "headline", None):
        _simple_text(slide, s.headline, left=ml, top=slide_h * 0.6, width=left_w, height=1.0, size=theme.size_headline, color=theme.ink_soft, font=theme.font_body, italic=style.headline_italic, line_spacing=1.1)

    _bullet_rows(
        slide, s.bullets,
        left=right_x, top=slide_h * 0.15, width=right_w, height=slide_h * 0.7,
        theme=theme, style=style, numbered=s.numbered,
    )

def _bullets_grid(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)
    
    n = len(s.bullets)
    cols = 2 if n <= 4 else 3
    rows = (n + cols - 1) // cols
    
    gap = 0.3
    content_w = slide_w - 2 * ml
    card_w = (content_w - gap * (cols - 1)) / cols
    avail_h = slide_h - ct - theme.margin_bottom - 0.3
    card_h = min(1.2, (avail_h - gap * (rows - 1)) / rows)
    
    for i, item in enumerate(s.bullets):
        r, c = divmod(i, cols)
        x = ml + c * (card_w + gap)
        y = ct + r * (card_h + gap)
        _draw_card(slide, x, y, card_w, card_h, theme, style)
        _rect(slide, MSO_SHAPE.RECTANGLE, x + 0.2, y + 0.2, 0.3, 0.04, fill=theme.accent)
        _simple_text(slide, item, left=x+0.2, top=y+0.35, width=card_w-0.4, height=card_h-0.4, size=theme.size_body-1, color=theme.ink, font=theme.font_body, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)


def render_bullets(prs, s: BulletsSlide, theme: Theme, style: Style, *, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    if style.layout_bullets == "magazine":
        _bullets_magazine(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    elif style.layout_bullets == "grid":
        _bullets_grid(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    else:
        _bullets_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    _set_notes(slide, s.notes)
    return slide


# ===================================================================
# MACRO-LAYOUT VARIANTS: COMPARISON
# ===================================================================

def _comparison_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    n = len(s.columns)
    content_w = slide_w - 2 * ml
    gap = 0.35
    col_w = (content_w - gap * (n - 1)) / n
    top = ct
    height = slide_h - top - theme.margin_bottom - 0.3

    for i, col in enumerate(s.columns):
        left = ml + i * (col_w + gap)
        _draw_card(slide, left, top, col_w, height, theme, style)
        header_color = theme.accent if i == 0 else theme.accent2
        _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, 0.62, fill=header_color)
        _simple_text(
            slide, col.label, left=left + 0.25, top=top, width=col_w - 0.5, height=0.62,
            size=theme.size_body + 2, color=theme.on_accent, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display,
        )
        _bullet_rows(
            slide, col.bullets, left=left + 0.28, top=top + 0.85, width=col_w - 0.56, height=height - 1.1,
            theme=theme, style=style, text_size=theme.size_body - 1, marker_color=header_color,
        )

def _comparison_split_vs(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _bg(slide, theme.accent)
    _rect(slide, MSO_SHAPE.RECTANGLE, slide_w/2, 0, slide_w/2, slide_h, fill=theme.background)
    
    if getattr(s, "kicker", None):
        _simple_text(slide, s.kicker.upper(), left=1.0, top=0.5, width=slide_w-2.0, height=0.3, size=theme.size_kicker, color=theme.on_accent_soft, bold=True, align=PP_ALIGN.CENTER, tracking=2.4)
    _simple_text(slide, s.title, left=1.0, top=0.8, width=slide_w-2.0, height=0.8, size=theme.size_title, color=theme.on_accent, bold=True, align=PP_ALIGN.CENTER, font=theme.font_display)
    
    n = len(s.columns)
    for i, col in enumerate(s.columns[:2]):
        is_left = (i == 0)
        fg_col = theme.on_accent if is_left else theme.ink
        marker_col = theme.on_accent_soft if is_left else theme.accent
        
        x = 0.5 if is_left else (slide_w/2 + 0.5)
        w = slide_w/2 - 1.0
        
        _simple_text(slide, col.label, left=x, top=2.0, width=w, height=0.8, size=theme.size_section, color=fg_col, bold=True, align=PP_ALIGN.CENTER, font=theme.font_display)
        _bullet_rows(slide, col.bullets, left=x+0.5, top=3.2, width=w-1.0, height=slide_h-4.0, theme=theme, style=style, text_size=theme.size_body, marker_color=marker_col)

    _rect(slide, MSO_SHAPE.OVAL, slide_w/2 - 0.4, slide_h/2 - 0.4, 0.8, 0.8, fill=theme.surface, line=theme.surface_edge)
    _simple_text(slide, "VS", left=slide_w/2 - 0.4, top=slide_h/2 - 0.4, width=0.8, height=0.8, size=16, color=theme.ink, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def _comparison_columns(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)
    
    n = len(s.columns)
    content_w = slide_w - 2 * ml
    gap = 0.6
    col_w = (content_w - gap * (n - 1)) / n
    top = ct
    height = slide_h - top - theme.margin_bottom - 0.3
    
    for i in range(1, n):
        x = ml + i * col_w + (i - 0.5) * gap
        _rect(slide, MSO_SHAPE.RECTANGLE, x, top + 0.2, 0.02, height - 0.4, fill=theme.surface_edge)
        
    for i, col in enumerate(s.columns):
        left = ml + i * (col_w + gap)
        header_color = theme.accent if i == 0 else theme.accent2
        _simple_text(
            slide, col.label, left=left, top=top, width=col_w, height=0.6,
            size=theme.size_body + 4, color=header_color, bold=True, anchor=MSO_ANCHOR.TOP, font=theme.font_display,
        )
        _bullet_rows(
            slide, col.bullets, left=left, top=top + 0.8, width=col_w, height=height - 0.8,
            theme=theme, style=style, text_size=theme.size_body, marker_color=header_color,
        )

def _comparison_rows(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)
    
    n = len(s.columns)
    top = ct
    avail_h = slide_h - top - theme.margin_bottom - 0.3
    gap = 0.2
    row_h = (avail_h - gap * (n - 1)) / n
    
    for i, col in enumerate(s.columns):
        y = top + i * (row_h + gap)
        header_color = theme.accent if i == 0 else theme.accent2
        _draw_card(slide, ml, y, slide_w - 2*ml, row_h, theme, style, fill_override=header_color)
        
        label_w = (slide_w - 2*ml) * 0.3
        _simple_text(slide, col.label, left=ml + 0.2, top=y, width=label_w - 0.4, height=row_h, size=theme.size_body + 4, color=theme.on_accent, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_display)
        
        t2 = Theme(name="tmp", ink=theme.on_accent, ink_soft=theme.on_accent_soft, background=theme.background, surface=theme.surface, surface_edge=theme.surface_edge, rule=theme.rule, accent=theme.on_accent_soft, accent_soft=theme.accent_soft, accent2=theme.accent2, on_accent=theme.on_accent, on_accent_soft=theme.on_accent_soft, code_bg=theme.code_bg, table_header_bg=theme.table_header_bg, table_zebra=theme.table_zebra, pygments_style=theme.pygments_style)
        _bullet_rows(slide, col.bullets, left=ml + label_w, top=y + 0.2, width=slide_w - 2*ml - label_w - 0.2, height=row_h - 0.4, theme=t2, style=style, text_size=theme.size_body-1)

def render_comparison(prs, s: ComparisonSlide, theme: Theme, style: Style, *, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    if style.layout_comparison == "split_vs" and len(s.columns) == 2:
        _comparison_split_vs(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    elif style.layout_comparison == "columns":
        _comparison_columns(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    elif style.layout_comparison == "rows":
        _comparison_rows(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    else:
        _comparison_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    _set_notes(slide, s.notes)
    return slide


# ===================================================================
# MACRO-LAYOUT VARIANTS: CODE
# ===================================================================

def _code_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    top = ct
    left = ml
    width = slide_w - 2 * ml
    height = slide_h - top - theme.margin_bottom - 0.3

    bg_hex = style_background_hex(theme.pygments_style)
    card_bg = RGBColor.from_string(bg_hex) if bg_hex else theme.code_bg
    _draw_card(slide, left, top, width, height, theme, style, fill_override=card_bg)

    if s.language and s.language != "text":
        pill = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + width - 1.4, top + 0.18, 1.2, 0.34, fill=theme.accent)
        ptf = pill.text_frame
        ptf.margin_left = Emu(0); ptf.margin_right = Emu(0); ptf.margin_top = Emu(0); ptf.margin_bottom = Emu(0)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        _run(pp, s.language.lower(), size=theme.size_caption, color=theme.on_accent, bold=True, font=theme.font_display)

    spans_per_line = highlight_lines(s.code, s.language, theme.pygments_style, str(theme.code_default))
    n = max(1, len(spans_per_line))
    pad = 0.28
    avail_h = height - 2 * pad
    line_h = min(0.34, avail_h / n)
    font_pt = min(theme.size_code, max(9, int(line_h * 72 * 0.62)))
    highlight = set(s.highlight_lines or [])

    text_left = left + pad + 0.15
    for idx, spans in enumerate(spans_per_line, start=1):
        ly = top + pad + (idx - 1) * line_h
        if idx in highlight:
            _rect(slide, MSO_SHAPE.RECTANGLE, left + 0.08, ly - 0.01, width - 0.16, line_h, fill=theme.code_line_highlight)
            _rect(slide, MSO_SHAPE.RECTANGLE, left + 0.08, ly - 0.01, 0.05, line_h, fill=theme.accent)
        box, tf = _textbox(slide, left=text_left, top=ly, width=width - pad - 0.3, height=line_h, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        if not spans:
            _run(p, " ", size=font_pt, color=card_bg, font=theme.font_mono)
        for sp in spans:
            _run(p, sp.text, size=font_pt, color=sp.color, bold=sp.bold, italic=sp.italic, font=theme.font_mono)

def _code_side_by_side(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _bg(slide, theme.background)
    
    left_w = slide_w * 0.35
    right_x = slide_w * 0.4
    right_w = slide_w * 0.6
    
    ml = _eff_margin(theme, style)
    if getattr(s, "kicker", None):
        _simple_text(slide, s.kicker.upper(), left=ml, top=1.0, width=left_w-ml, height=0.4, size=theme.size_kicker, color=theme.accent, bold=True, tracking=2.4)
    _simple_text(slide, s.title, left=ml, top=1.4, width=left_w-ml, height=1.6, size=theme.size_title, color=theme.ink, bold=True, font=theme.font_display)
    if getattr(s, "headline", None):
        _simple_text(slide, s.headline, left=ml, top=3.0, width=left_w-ml, height=3.0, size=theme.size_headline, color=theme.ink_soft, font=theme.font_body, line_spacing=1.1)
        
    bg_hex = style_background_hex(theme.pygments_style)
    card_bg = RGBColor.from_string(bg_hex) if bg_hex else theme.code_bg
    _rect(slide, MSO_SHAPE.RECTANGLE, right_x, 0, right_w, slide_h, fill=card_bg)
    
    spans_per_line = highlight_lines(s.code, s.language, theme.pygments_style, str(theme.code_default))
    n = max(1, len(spans_per_line))
    pad = 0.5
    avail_h = slide_h - 2 * pad
    line_h = min(0.34, avail_h / n)
    font_pt = min(theme.size_code + 2, max(9, int(line_h * 72 * 0.62)))
    highlight = set(s.highlight_lines or [])

    text_left = right_x + pad
    for idx, spans in enumerate(spans_per_line, start=1):
        ly = pad + (idx - 1) * line_h
        if idx in highlight:
            _rect(slide, MSO_SHAPE.RECTANGLE, right_x, ly - 0.01, right_w, line_h, fill=theme.code_line_highlight)
            _rect(slide, MSO_SHAPE.RECTANGLE, right_x, ly - 0.01, 0.08, line_h, fill=theme.accent)
        box, tf = _textbox(slide, left=text_left, top=ly, width=right_w - pad - 0.2, height=line_h, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        if not spans:
            _run(p, " ", size=font_pt, color=card_bg, font=theme.font_mono)
        for sp in spans:
            _run(p, sp.text, size=font_pt, color=sp.color, bold=sp.bold, italic=sp.italic, font=theme.font_mono)

def render_code(prs, s: CodeSlide, theme: Theme, style: Style, *, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    if style.layout_code == "side_by_side":
        _code_side_by_side(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    else:
        _code_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    _set_notes(slide, s.notes)
    return slide


def render_diagram(prs, s: DiagramSlide, theme: Theme, style: Style, *, base_dir, mermaid_renderer=None, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    top = ct
    box_w = slide_w - 2 * ml
    box_h = slide_h - top - theme.margin_bottom - 0.3
    box = (ml, top, box_w, box_h)

    if mermaid_renderer is None:
        _placeholder(slide, box, theme, "[diagram rendering disabled \u2014 install mermaid]")
        _set_notes(slide, s.notes)
        return slide

    with tempfile.TemporaryDirectory() as td:
        out = Path(td) / "diagram.png"
        try:
            mermaid_renderer(s.mermaid, out)
            _place_image_contain(slide, out, box, theme, shadow=False)
        except Exception as e:  # noqa: BLE001
            _placeholder(slide, box, theme, f"[diagram failed: {e}]", color=RGBColor.from_string("C0392B"))
    _set_notes(slide, s.notes)
    return slide


# ===================================================================
# MACRO-LAYOUT VARIANTS: IMAGE
# ===================================================================

def _image_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    img_path = Path(s.path)
    if not img_path.is_absolute():
        img_path = (base_dir / s.path).resolve()

    top = ct
    full_w = slide_w - 2 * ml
    full_h = slide_h - top - theme.margin_bottom - 0.3
    cap_h = 0.4 if s.caption else 0.0

    layout = s.layout
    if layout == "auto":
        layout = "right" if s.bullets else "full"

    if layout in ("right", "left") and s.bullets:
        gap = 0.4
        col_w = (full_w - gap) / 2
        if layout == "right":
            text_x, img_x = ml, ml + col_w + gap
        else:
            img_x, text_x = ml, ml + col_w + gap
        _bullet_rows(
            slide, s.bullets, left=text_x, top=top + 0.1, width=col_w, height=full_h - 0.2,
            theme=theme, style=style,
        )
        img_box = (img_x, top, col_w, full_h - cap_h)
        if img_path.exists():
            _place_image_contain(slide, img_path, img_box, theme)
        else:
            _placeholder(slide, img_box, theme, f"[missing image: {s.path}]", color=RGBColor.from_string("C0392B"))
        if s.caption:
            _simple_text(slide, s.caption, left=img_x, top=top + full_h - cap_h + 0.05, width=col_w, height=cap_h, size=theme.size_caption, color=theme.ink_soft, align=PP_ALIGN.CENTER, font=theme.font_body)
    else:
        img_box = (ml, top, full_w, full_h - cap_h)
        if img_path.exists():
            _place_image_contain(slide, img_path, img_box, theme)
        else:
            _placeholder(slide, img_box, theme, f"[missing image: {s.path}]", color=RGBColor.from_string("C0392B"))
        if s.caption:
            _simple_text(slide, s.caption, left=ml, top=top + full_h - cap_h + 0.05, width=full_w, height=cap_h, size=theme.size_caption, color=theme.ink_soft, align=PP_ALIGN.CENTER, font=theme.font_body)

def _image_split_bleed(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir):
    _bg(slide, theme.background)
    
    img_path = Path(s.path)
    if not img_path.is_absolute():
        img_path = (base_dir / s.path).resolve()
        
    left_w = slide_w * 0.5
    right_x = slide_w * 0.55
    right_w = slide_w * 0.4
    
    if img_path.exists():
        _place_image_contain(slide, img_path, (0, 0, left_w, slide_h), theme, shadow=False)
    else:
        _placeholder(slide, (0, 0, left_w, slide_h), theme, f"[missing image: {s.path}]", color=RGBColor.from_string("C0392B"))

    _simple_text(slide, s.title, left=right_x, top=slide_h*0.2, width=right_w, height=1.0, size=theme.size_title, color=theme.ink, bold=True, font=theme.font_display)
    if getattr(s, "headline", None):
        _simple_text(slide, s.headline, left=right_x, top=slide_h*0.35, width=right_w, height=0.6, size=theme.size_headline, color=theme.ink_soft, font=theme.font_body)
        
    if s.bullets:
        _bullet_rows(slide, s.bullets, left=right_x, top=slide_h*0.5, width=right_w, height=slide_h*0.4, theme=theme, style=style)

def _image_immersive(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir):
    _bg(slide, theme.background)
    img_path = Path(s.path)
    if not img_path.is_absolute():
        img_path = (base_dir / s.path).resolve()
        
    if img_path.exists():
        _place_image_contain(slide, img_path, (0, 0, slide_w, slide_h), theme, shadow=False)
    else:
        _placeholder(slide, (0, 0, slide_w, slide_h), theme, f"[missing image: {s.path}]", color=RGBColor.from_string("C0392B"))
        
    # Dark overlay
    try:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(slide_w), Inches(slide_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0,0,0)
        shp.fill.fore_color.theme_color = 1 # hack to make it solid
        shp.line.fill.background()
    except Exception:
        pass
        
    card_w = slide_w * 0.7
    card_h = slide_h * 0.7
    cx = (slide_w - card_w) / 2
    cy = (slide_h - card_h) / 2
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h, fill=RGBColor(20,20,30), line=RGBColor(80,80,100))
    
    _simple_text(slide, s.title, left=cx+0.5, top=cy+0.5, width=card_w-1.0, height=1.0, size=theme.size_title, color=RGBColor(255,255,255), bold=True, font=theme.font_display, align=PP_ALIGN.CENTER)
    if s.bullets:
        t2 = Theme(name="tmp", ink=RGBColor(255,255,255), ink_soft=RGBColor(200,200,200), background=theme.background, surface=theme.surface, surface_edge=theme.surface_edge, rule=theme.rule, accent=theme.accent, accent_soft=theme.accent_soft, accent2=theme.accent2, on_accent=theme.on_accent, on_accent_soft=theme.on_accent_soft, code_bg=theme.code_bg, table_header_bg=theme.table_header_bg, table_zebra=theme.table_zebra, pygments_style=theme.pygments_style)
        _bullet_rows(slide, s.bullets, left=cx+1.0, top=cy+2.0, width=card_w-2.0, height=card_h-2.5, theme=t2, style=style)

def render_image(prs, s: ImageSlide, theme: Theme, style: Style, *, base_dir, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    if style.layout_image == "split_bleed":
        _image_split_bleed(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir)
    elif style.layout_image == "immersive":
        _image_immersive(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir)
    else:
        _image_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total, base_dir)
    _set_notes(slide, s.notes, source=s.source)
    return slide


# ===================================================================
# MACRO-LAYOUT VARIANTS: TABLE
# ===================================================================

def _table_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    top = ct
    width = slide_w - 2 * ml
    n_rows = len(s.rows) + 1
    row_h = min(0.55, (slide_h - top - theme.margin_bottom - 0.3) / n_rows)
    height = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, len(s.columns), Inches(ml), Inches(top), Inches(width), Inches(height))
    table = shape.table
    table.first_row = False; table.horz_banding = False

    for j, name in enumerate(s.columns):
        c = table.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = theme.table_header_bg
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.12); c.margin_right = Inches(0.12)
        _run(c.text_frame.paragraphs[0], str(name), size=theme.size_body - 1, color=theme.table_header_fg, bold=True, font=theme.font_display)

    for i, row in enumerate(s.rows, start=1):
        zebra = theme.table_zebra if i % 2 == 1 else theme.background
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = zebra
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.12)
            _run(c.text_frame.paragraphs[0], str(val), size=theme.size_body - 2, color=theme.ink, font=theme.font_body)

def _table_clean(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    top = ct
    width = slide_w - 2 * ml
    n_rows = len(s.rows) + 1
    row_h = min(0.6, (slide_h - top - theme.margin_bottom - 0.3) / n_rows)
    height = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, len(s.columns), Inches(ml), Inches(top), Inches(width), Inches(height))
    table = shape.table
    table.first_row = False; table.horz_banding = False

    for j, name in enumerate(s.columns):
        c = table.cell(0, j)
        c.fill.background()
        c.vertical_anchor = MSO_ANCHOR.BOTTOM
        c.margin_left = Inches(0.12); c.margin_right = Inches(0.12)
        _run(c.text_frame.paragraphs[0], str(name), size=theme.size_body - 1, color=theme.ink_soft, bold=True, font=theme.font_display)

    for i, row in enumerate(s.rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.fill.background()
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.12)
            _run(c.text_frame.paragraphs[0], str(val), size=theme.size_body - 1, color=theme.ink, font=theme.font_body)

    _rect(slide, MSO_SHAPE.RECTANGLE, ml, top + row_h - 0.05, width, 0.02, fill=theme.ink_soft) 
    for i in range(1, n_rows):
        _rect(slide, MSO_SHAPE.RECTANGLE, ml, top + (i+1)*row_h - 0.02, width, 0.01, fill=theme.rule)

def _table_floating_cards(slide, s, theme, style, slide_w, slide_h, footer, page, total):
    _content_chrome(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)

    top = ct
    width = slide_w - 2 * ml
    n_rows = len(s.rows)
    n_cols = max(1, len(s.columns))
    col_w = width / n_cols
    
    gap = 0.2
    avail_h = slide_h - top - theme.margin_bottom - 0.3
    row_h = min(0.8, (avail_h - gap * (n_rows - 1)) / n_rows)
    
    for j, name in enumerate(s.columns):
        _simple_text(slide, str(name), left=ml + j*col_w, top=top, width=col_w, height=0.3, size=theme.size_caption, color=theme.accent, bold=True, anchor=MSO_ANCHOR.BOTTOM, font=theme.font_display)
        
    top += 0.4
    
    for i, row in enumerate(s.rows):
        y = top + i * (row_h + gap)
        _draw_card(slide, ml, y, width, row_h, theme, style)
        for j, val in enumerate(row):
            _simple_text(slide, str(val), left=ml + j*col_w + 0.15, top=y, width=col_w - 0.3, height=row_h, size=theme.size_body - 1, color=theme.ink, anchor=MSO_ANCHOR.MIDDLE, font=theme.font_body)

def render_table(prs, s: TableSlide, theme: Theme, style: Style, *, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    if style.layout_table == "clean":
        _table_clean(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    elif style.layout_table == "floating_cards":
        _table_floating_cards(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    else:
        _table_classic(slide, s, theme, style, slide_w, slide_h, footer, page, total)
    _set_notes(slide, s.notes)
    return slide


def render_quote(prs, s: QuoteSlide, theme: Theme, style: Style, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    fn = _QUOTE_VARIANTS.get(style.quote_variant, _quote_big_mark)
    fn(slide, s, theme, style, slide_w, slide_h)
    _set_notes(slide, s.notes)
    return slide


def render_summary(prs, s: SummarySlide, theme: Theme, style: Style, *, footer=None, page=None, total=None, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    ml = _eff_margin(theme, style)
    ct = _eff_content_top(style)
    _bg(slide, theme.background)
    _kicker(slide, "Key takeaways", theme, style, slide_w)
    _title_chrome(slide, s.title, theme, style, slide_w)
    _footer(slide, theme, style, slide_w, slide_h, footer, page, total)

    top = ct
    width = slide_w - 2 * ml
    avail = slide_h - top - theme.margin_bottom - 0.3
    n = len(s.points)
    row_h = min(1.2, avail / n)
    for i, point in enumerate(s.points):
        ry = top + i * row_h
        _simple_text(
            slide, f"{i+1:02d}",
            left=ml, top=ry,
            width=1.0, height=row_h,
            size=theme.size_section, color=theme.accent, bold=True, font=theme.font_display,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _simple_text(
            slide, point,
            left=ml + 1.15, top=ry,
            width=width - 1.15, height=row_h,
            size=theme.size_body_lg, color=theme.ink, font=theme.font_body,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02,
        )
        if i < n - 1:
            _rect(slide, MSO_SHAPE.RECTANGLE, ml + 1.15, ry + row_h - 0.02, width - 1.15, 0.01, fill=theme.rule)
    _set_notes(slide, s.notes)
    return slide


def render_qa(prs, s: QASlide, theme: Theme, style: Style, **_):
    slide = _new_slide(prs)
    slide_w, slide_h = _slide_size_in(prs)
    _bg(slide, theme.accent)
    _simple_text(
        slide, s.title,
        left=theme.margin_left, top=slide_h * 0.34,
        width=slide_w - 2 * theme.margin_left, height=1.6,
        size=theme.size_qa, color=theme.on_accent, bold=True,
        align=PP_ALIGN.CENTER, font=theme.font_display, anchor=MSO_ANCHOR.MIDDLE,
    )
    if s.contact:
        pill = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, slide_w / 2 - 2.0, slide_h * 0.62, 4.0, 0.5, fill=theme.on_accent)
        ptf = pill.text_frame
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        _run(pp, s.contact, size=theme.size_attribution, color=theme.accent, bold=True, font=theme.font_body)
    _set_notes(slide, s.notes)
    return slide


# ---------------------------------------------------------------------------
# dispatch
# ---------------------------------------------------------------------------


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "statement": render_statement,
    "bullets": render_bullets,
    "comparison": render_comparison,
    "code": render_code,
    "diagram": render_diagram,
    "image": render_image,
    "table": render_table,
    "quote": render_quote,
    "summary": render_summary,
    "qa": render_qa,
}


import dataclasses
from pptx.enum.shapes import PP_PLACEHOLDER

def render_slide(
    prs,
    slide_model: Slide,
    theme: Theme,
    style: Style,
    *,
    base_dir: Path,
    mermaid_renderer=None,
    footer=None,
    page=None,
    total=None,
    section_index=None,
    meta_config=None,
):
    if meta_config:
        layout_idx = meta_config['mappings'].get(slide_model.type, 0)
        sa = meta_config['safe_areas'].get(layout_idx)
        
        # Simple Slides: Use native placeholders, completely bypass custom drawing
        if slide_model.type in ("title", "section", "statement", "bullets", "quote", "summary", "qa"):
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Find and fill title
            for shape in slide.placeholders:
                if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    shape.text = getattr(slide_model, 'title', '') or getattr(slide_model, 'text', '') or getattr(slide_model, 'quote', '')
                    break
            
            # Find and fill body
            if slide_model.type == "bullets" and hasattr(slide_model, 'bullets'):
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                        tf = shape.text_frame
                        tf.text = ""
                        for b in slide_model.bullets:
                            p = tf.add_paragraph()
                            p.text = str(b)
                            p.level = 0
                        break
            elif slide_model.type == "qa" and hasattr(slide_model, 'answer'):
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                        shape.text = slide_model.answer
                        break

            _set_notes(slide, slide_model.notes)
            return slide
            
        # Complex Slides: Use safe area constraints
        # Create a shallow copy of theme and style using dataclasses.replace
        theme = dataclasses.replace(theme)
        style = dataclasses.replace(style)
        
        if sa:
            slide_w, slide_h = _slide_size_in(prs)
            theme = dataclasses.replace(
                theme,
                margin_left=sa['left'],
                margin_right=slide_w - (sa['left'] + sa['width']),
                margin_bottom=slide_h - (sa['top'] + sa['height'])
            )
            style = dataclasses.replace(
                style,
                content_top_offset=sa['top'] - CONTENT_TOP,
                margin_scale=1.0
            )

        # We need to override pick_blank_layout just for this run so it uses our mapped layout
        original_pick_blank = globals().get('pick_blank_layout')
        globals()['pick_blank_layout'] = lambda p: p.slide_layouts[layout_idx]

        try:
            fn = RENDERERS[slide_model.type]
            res = fn(
                prs,
                slide_model,
                theme,
                style,
                base_dir=base_dir,
                mermaid_renderer=mermaid_renderer,
                footer=footer,
                page=page,
                total=total,
                section_index=section_index,
            )
            
            # Since we used the layout directly, we need to hide native placeholders to avoid collision
            # Our custom renderers draw their own title and text boxes.
            for shape in res.placeholders:
                shape.element.getparent().remove(shape.element)
                
            return res
        finally:
            if original_pick_blank:
                globals()['pick_blank_layout'] = original_pick_blank


    fn = RENDERERS[slide_model.type]
    return fn(
        prs,
        slide_model,
        theme,
        style,
        base_dir=base_dir,
        mermaid_renderer=mermaid_renderer,
        footer=footer,
        page=page,
        total=total,
        section_index=section_index,
    )
