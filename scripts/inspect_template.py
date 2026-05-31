"""Inspect a .pptx template and dump its layouts/placeholders/theme as YAML.

Run this before drafting outline.yaml when the user supplies a company template,
so the agent knows which layouts exist and which slide types can map to them.

Usage:
    python scripts/inspect_template.py path/to/company-template.pptx
    python scripts/inspect_template.py path/to/company-template.pptx -o spec.yaml
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.util import Emu


def _emu_to_in(v: int | None) -> float | None:
    return None if v is None else round(Emu(v).inches, 3)


def _placeholder_info(ph) -> dict[str, Any]:
    return {
        "idx": ph.placeholder_format.idx,
        "type": str(ph.placeholder_format.type),
        "name": ph.name,
        "left_in": _emu_to_in(ph.left),
        "top_in": _emu_to_in(ph.top),
        "width_in": _emu_to_in(ph.width),
        "height_in": _emu_to_in(ph.height),
    }


def _layout_info(layout) -> dict[str, Any]:
    return {
        "name": layout.name,
        "placeholders": [_placeholder_info(p) for p in layout.placeholders],
    }


def inspect(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    return {
        "file": str(path),
        "slide_size_in": {
            "width": _emu_to_in(prs.slide_width),
            "height": _emu_to_in(prs.slide_height),
        },
        "layout_count": len(prs.slide_layouts),
        "layouts": [
            {"index": i, **_layout_info(layout)}
            for i, layout in enumerate(prs.slide_layouts)
        ],
        "suggested_mapping": _suggest_mapping(prs),
    }


# Heuristic mapping from slide type -> layout name fragments, in priority order.
# build_deck.py uses these same hints via lib.layouts; this file just exposes
# the best guesses to the human reader.
_MAPPING_HINTS: dict[str, list[str]] = {
    "title": ["title slide", "title"],
    "section": ["section header", "section"],
    "statement": ["title only", "centered", "big idea"],
    "bullets": ["title and content", "content"],
    "comparison": ["comparison", "two content"],
    "code": ["title and content", "content"],
    "diagram": ["title and content", "picture", "content"],
    "image": ["picture with caption", "picture", "title and content"],
    "table": ["title and content", "content"],
    "quote": ["title only", "blank"],
    "summary": ["title and content", "content"],
    "qa": ["title only", "section header"],
}


def _suggest_mapping(prs) -> dict[str, dict[str, Any]]:
    layouts_lower = [(i, layout.name.lower()) for i, layout in enumerate(prs.slide_layouts)]
    result: dict[str, dict[str, Any]] = {}
    for slide_type, fragments in _MAPPING_HINTS.items():
        match_idx, match_name, match_frag = None, None, None
        for frag in fragments:
            for i, name in layouts_lower:
                if frag in name:
                    match_idx, match_name, match_frag = i, prs.slide_layouts[i].name, frag
                    break
            if match_idx is not None:
                break
        result[slide_type] = {
            "layout_index": match_idx,
            "layout_name": match_name,
            "matched_on": match_frag,
        }
    return result


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Inspect a .pptx template")
    parser.add_argument("template", type=Path)
    parser.add_argument(
        "-o", "--output", type=Path, help="Write YAML here (default: stdout)"
    )
    args = parser.parse_args(argv)

    if not args.template.exists():
        print(f"error: {args.template} not found", file=sys.stderr)
        return 2

    spec = inspect(args.template)
    text = yaml.safe_dump(spec, sort_keys=False, allow_unicode=True)

    if args.output:
        args.output.write_text(text, encoding="utf-8")
        print(f"wrote {args.output}")
    else:
        sys.stdout.write(text)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
