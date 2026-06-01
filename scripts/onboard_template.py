"""Onboard a .pptx template: map layouts, compute safe areas, and save a .meta.yaml.

Usage:
    python scripts/onboard_template.py path/to/template.pptx
    python scripts/onboard_template.py path/to/template.pptx --auto
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu

_MAPPING_HINTS: dict[str, list[str]] = {
    "title": ["title slide", "title"],
    "section": ["section header", "section"],
    "statement": ["title only", "centered", "big idea", "blank"],
    "bullets": ["title and content", "content"],
    "comparison": ["comparison", "two content"],
    "code": ["title and content", "content"],
    "diagram": ["title and content", "picture", "content"],
    "image": ["picture with caption", "picture", "title and content"],
    "table": ["title and content", "content"],
    "quote": ["title only", "blank"],
    "summary": ["title and content", "content"],
    "qa": ["title only", "section header", "blank"],
}

def _emu_to_in(v: int | None) -> float | None:
    return None if v is None else round(Emu(v).inches, 3)

def _score_layout(layout_idx: int, layout, slide_type: str) -> int:
    score = 0
    name = layout.name.lower()

    # 1. Tally placeholders
    title_count = 0
    body_count = 0
    for shape in layout.placeholders:
        if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title_count += 1
        elif shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            body_count += 1

    # 2. Semantic Hinting (Name Matching)
    if slide_type in _MAPPING_HINTS:
        for frag in _MAPPING_HINTS[slide_type]:
            if frag in name:
                score += 50
                break

    # 3. Structural Signature Scoring
    if slide_type == "title":
        if title_count >= 1: score += 20
        if body_count == 0: score += 30
        if body_count > 0: score -= 50
        if layout_idx == 0: score += 50  # Layout 0 is almost always the title slide
    elif slide_type == "comparison":
        if title_count >= 1: score += 10
        if body_count >= 2: score += 50
        elif body_count == 1: score -= 20
        else: score -= 50
        if layout_idx == 0: score -= 50
    elif slide_type == "section":
        if title_count >= 1: score += 20
        if body_count == 0: score += 30
        if body_count > 0: score -= 50
        if "section" in name or "transition" in name: score += 30
        if layout_idx == 0: score -= 50
    else:
        # bullets, code, statement, diagram, image, table, summary, qa
        if title_count >= 1: score += 10
        if body_count == 1: score += 40
        elif body_count > 1: score -= 30  # Penalty for too many bodies
        else: score -= 30 # Penalty for no body
        if layout_idx == 0: score -= 5000  # NEVER use the title slide for content

    return score

def _suggest_mappings(prs) -> dict[str, int]:
    used_layout_indices = {prs.slide_layouts.index(slide.slide_layout) for slide in prs.slides}

    mappings: dict[str, int] = {}
    for slide_type in _MAPPING_HINTS.keys():
        best_idx = 0
        best_score = -9999
        
        for i, layout in enumerate(prs.slide_layouts):
            s = _score_layout(i, layout, slide_type)
            if i in used_layout_indices:
                s += 1000  # HUGE boost for layouts actually used in the template
                
            # Give a slight edge to earlier layouts if scores are tied
            if s > best_score:
                best_score = s
                best_idx = i
                
        mappings[slide_type] = best_idx
    return mappings

def _intersect(box1, box2):
    l1, t1, r1, b1 = box1
    l2, t2, r2, b2 = box2
    return max(l1, l2) < min(r1, r2) and max(t1, t2) < min(b1, b2)

def _calc_safe_area(prs, layout_idx: int) -> dict[str, float]:
    layout = prs.slide_layouts[layout_idx]
    sw = _emu_to_in(prs.slide_width)
    sh = _emu_to_in(prs.slide_height)

    # If there are body placeholders, use the one with the largest area as the safe area
    best_body_ph = None
    best_area = -1
    for ph in layout.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            area = _emu_to_in(ph.width) * _emu_to_in(ph.height)
            if area > best_area:
                best_area = area
                best_body_ph = ph

    if best_body_ph:
        return {
            "left": _emu_to_in(best_body_ph.left),
            "top": _emu_to_in(best_body_ph.top),
            "width": _emu_to_in(best_body_ph.width),
            "height": _emu_to_in(best_body_ph.height)
        }

    # Otherwise, compute bounding box using heuristics
    margin = 0.5
    safe_left, safe_top, safe_right, safe_bottom = margin, margin, sw - margin, sh - margin

    def apply_avoidance(shape, is_ph):
        nonlocal safe_left, safe_top, safe_right, safe_bottom
        sl, st = _emu_to_in(shape.left), _emu_to_in(shape.top)
        sw2, sh2 = _emu_to_in(shape.width), _emu_to_in(shape.height)
        sr, sb = sl + sw2, st + sh2

        if is_ph:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                # Title typically at top
                safe_top = max(safe_top, sb + margin)
            elif ph_type in (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE):
                # Footers typically at bottom
                safe_bottom = min(safe_bottom, st - margin)
        else:
            # Static decoration shape
            area = sw2 * sh2
            slide_area = sw * sh
            if area > slide_area * 0.5:
                return # Ignore huge background images

            # If it intersects our safe area, push the closest boundary inward
            if _intersect((safe_left, safe_top, safe_right, safe_bottom), (sl, st, sr, sb)):
                d_left = abs(safe_left - sr)
                d_top = abs(safe_top - sb)
                d_right = abs(safe_right - sl)
                d_bottom = abs(safe_bottom - st)

                min_d = min(d_left, d_top, d_right, d_bottom)
                if min_d == d_left: safe_left = sr + 0.1
                elif min_d == d_right: safe_right = sl - 0.1
                elif min_d == d_top: safe_top = sb + 0.1
                elif min_d == d_bottom: safe_bottom = st - 0.1

    # Apply avoidance for master static shapes
    for shape in layout.slide_master.shapes:
        apply_avoidance(shape, shape.is_placeholder)
    
    # Apply avoidance for layout shapes
    for shape in layout.shapes:
        apply_avoidance(shape, shape.is_placeholder)

    # Ensure valid bounds
    if safe_right <= safe_left: safe_right = safe_left + 1.0
    if safe_bottom <= safe_top: safe_bottom = safe_top + 1.0

    return {
        "left": round(safe_left, 3),
        "top": round(safe_top, 3),
        "width": round(safe_right - safe_left, 3),
        "height": round(safe_bottom - safe_top, 3)
    }

def main(argv: list[str] | None = None) -> int:
    import sys
    if sys.stdout.encoding.lower() != 'utf-8':
        try:
            sys.stdout.reconfigure(encoding='utf-8')
        except AttributeError:
            pass

    parser = argparse.ArgumentParser(description="Onboard a .pptx template interactively or headlessly.")
    parser.add_argument("template", type=Path)
    parser.add_argument("--auto", action="store_true", help="Run headlessly without interactive prompts.")
    args = parser.parse_args(argv)

    if not args.template.exists():
        print(f"error: {args.template} not found", file=sys.stderr)
        return 1

    print(f"Loading template: {args.template}")
    prs = Presentation(str(args.template))
    
    print(f"Found {len(prs.slide_layouts)} layouts.")
    for i, l in enumerate(prs.slide_layouts):
        print(f"  [{i}] {l.name}")

    mappings = _suggest_mappings(prs)
    
    if not args.auto:
        print("\n--- Mapping Slide Types to Layouts ---")
        for slide_type in list(mappings.keys()):
            suggested = mappings[slide_type]
            layout_name = prs.slide_layouts[suggested].name
            ans = input(f"Map '{slide_type}' to [{suggested}] {layout_name}? [Y/n/index] ").strip()
            if ans.lower() == 'n':
                idx = input(f"  Enter layout index for '{slide_type}': ").strip()
                if idx.isdigit():
                    mappings[slide_type] = int(idx)
            elif ans.isdigit():
                mappings[slide_type] = int(ans)
    
    print("\n--- Computing Safe Areas ---")
    safe_areas = {}
    unique_layouts_used = set(mappings.values())
    for idx in unique_layouts_used:
        layout_name = prs.slide_layouts[idx].name
        sa = _calc_safe_area(prs, idx)
        safe_areas[idx] = sa
        if args.auto:
            print(f"  [{idx}] {layout_name}: {sa}")

    meta = {
        "template_file": args.template.name,
        "slide_size_in": {
            "width": _emu_to_in(prs.slide_width),
            "height": _emu_to_in(prs.slide_height)
        },
        "mappings": mappings,
        "safe_areas": safe_areas
    }

    out_file = args.template.with_suffix(".meta.yaml")
    with open(out_file, "w", encoding="utf-8") as f:
        yaml.safe_dump(meta, f, sort_keys=False)
    
    print(f"\nSuccess! Template metadata saved to: {out_file}")
    return 0

if __name__ == "__main__":
    sys.exit(main())
